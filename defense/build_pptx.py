"""
LMSIS Dissertation Defense — .pptx Generator
Generates a widescreen 16:9 (13.33in x 7.5in) PowerPoint using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import math
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
BG       = RGBColor(0x05, 0x08, 0x10)
PANEL    = RGBColor(0x0C, 0x11, 0x1E)
BORDER   = RGBColor(0x1C, 0x29, 0x40)
TEXT     = RGBColor(0xEE, 0xF2, 0xFF)
MUTED    = RGBColor(0x4A, 0x63, 0x80)
TEAL     = RGBColor(0x00, 0xC4, 0x7D)
AMBER    = RGBColor(0xF5, 0xA6, 0x23)
BLUE     = RGBColor(0x3D, 0x8E, 0xF8)
RED      = RGBColor(0xE8, 0x39, 0x4A)
GOLD     = RGBColor(0xC8, 0xA8, 0x4B)

# Slide size: 13.33in x 7.5in (16:9 widescreen)
W_IN, H_IN = 13.33, 7.5
W = Inches(W_IN)
H = Inches(H_IN)

MARGIN = Inches(W_IN * 0.09)   # 9% margin

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


# ── Low-level shape helpers ────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Inter", font_size=24, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, wrap=True,
                italic=False, mono=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "JetBrains Mono" if mono else font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_name="Inter", font_size=24, bold=False,
             color=TEXT, align=PP_ALIGN.LEFT, italic=False, mono=False, space_before=None):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = "JetBrains Mono" if mono else font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_footer_num(slide, num, total=15):
    """Tiny slide-number footer, bottom-right."""
    label = f"{num} / {total}"
    add_textbox(slide, label,
                W - Inches(1.2), H - Inches(0.35),
                Inches(1.0), Inches(0.28),
                font_size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_glow_rect(slide, left, top, width, height, color: RGBColor, alpha_stops=None):
    """Simulate a soft glow behind a hero number with a gradient rect."""
    # Outer soft box  
    glow = add_rect(slide, left - Inches(0.3), top - Inches(0.15),
                    width + Inches(0.6), height + Inches(0.3),
                    fill_color=PANEL, line_color=color, line_width_pt=1.5)
    return glow


# ── Map-pin SVG helper (embedded as EMF-less inline) ─────────────────────────
def add_mappin(slide, cx_in, cy_in, size_in=0.45, color: RGBColor = TEAL):
    """
    Draw a map-pin shape (circle + triangle pointer) at (cx_in, cy_in) inches.
    We use two overlapping pptx shapes.
    """
    r = Inches(size_in * 0.5)
    left = Inches(cx_in) - r
    top  = Inches(cy_in) - r
    # head circle
    circle = slide.shapes.add_shape(9,  # oval
                                     left, top, r * 2, r * 2)
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # tail triangle pointing down
    tri = slide.shapes.add_shape(5,   # isoceles triangle
                                  left + r - Inches(size_in * 0.22),
                                  top  + r * 1.5,
                                  Inches(size_in * 0.44),
                                  Inches(size_in * 0.55))
    # rotate 180 so point faces down
    sp_el = tri._element
    spPr = sp_el.find(qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', str(int(math.pi * 10800000)))  # 180 deg in 60000ths/deg  = 10800000 * 180 / 180
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    return circle, tri


def add_checkmark(slide, cx_in, cy_in, size_in=0.25, color=TEAL):
    """Draw a simple check-mark tick using a line shape approximation."""
    # Use a rounded rectangle with a ✓ character
    box = add_rect(slide,
                   Inches(cx_in - size_in / 2), Inches(cy_in - size_in / 2),
                   Inches(size_in), Inches(size_in),
                   fill_color=color)
    box.line.fill.background()
    tf = box.text_frame
    tf.paragraphs[0].text = "✓"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(size_in * 72 * 0.65)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x05, 0x08, 0x10)
    return box


# ══════════════════════════════════════════════════════════════════════════════
# Build slides
# ══════════════════════════════════════════════════════════════════════════════

prs = new_prs()
blank_layout = prs.slide_layouts[6]   # completely blank layout


# ────────────────────────────────────────────────────────── SLIDE 1 — Title ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)

# Subtle horizontal teal rule to anchor the title block
rule = add_rect(slide,
                MARGIN, Inches(2.1),
                W - 2 * MARGIN, Pt(2),
                fill_color=TEAL)

# Main title
tb = slide.shapes.add_textbox(MARGIN, Inches(1.0), W - 2 * MARGIN, Inches(1.3))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Predictive Risk Intelligence for Metabolic Screening in Diabetes"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True
r.font.color.rgb = TEXT

# Subtitle LMSIS
tb2 = slide.shapes.add_textbox(MARGIN, Inches(2.35), W - 2 * MARGIN, Inches(0.65))
tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "LMSIS — Latent Metabolic State Inference System"
r2.font.name = "Inter"; r2.font.size = Pt(26); r2.font.bold = False
r2.font.color.rgb = TEAL

# Decorative teal/amber colour bar row beneath subtitle
for i, col in enumerate([TEAL, AMBER, BLUE, RED]):
    add_rect(slide,
             MARGIN + (W - 2 * MARGIN) * i / 4,
             Inches(3.05),
             (W - 2 * MARGIN) / 4,
             Pt(3),
             fill_color=col)

# Presenter placeholder block
ph_top = Inches(4.5)
ph_h   = Inches(0.55)
for placeholder, idx in [("[Presenter Name]", 0), ("[Date]", 1)]:
    tb_ph = slide.shapes.add_textbox(
        Inches(W_IN * 0.3), ph_top + Inches(idx * 0.7),
        Inches(W_IN * 0.4), ph_h)
    tf_ph = tb_ph.text_frame
    p_ph = tf_ph.paragraphs[0]
    p_ph.alignment = PP_ALIGN.CENTER
    r_ph = p_ph.add_run()
    r_ph.text = placeholder
    r_ph.font.name = "Inter"; r_ph.font.size = Pt(22)
    r_ph.font.color.rgb = MUTED; r_ph.font.italic = True

# map-pin motif (small, teal) — bottom-left decorative
add_mappin(slide, 1.1, 6.8, size_in=0.38, color=TEAL)


# ────────────────────────────────────────────────────── SLIDE 2 — The Paradox ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 2)

lines = [
    ("Normal BMI. Normal cholesterol. No symptoms.", TEXT,   False, 34),
    ("Their doctor sent them home.",                 MUTED,  False, 32),
    ("They had severe liver disease.",               RED,    True,  38),
]
start_y = Inches(2.3)
for txt, col, bold, sz in lines:
    tb = slide.shapes.add_textbox(MARGIN, start_y, W - 2 * MARGIN, Inches(0.85))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = txt; r.font.name = "Inter"; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.color.rgb = col
    start_y += Inches(1.1)


# ──────────────────────────────────────────────── SLIDE 3 — Why BMI Misses ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 3)

# Title
tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Why BMI Misses This"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

# Two body silhouettes: just two rounded rects with labels
fig_w = Inches(1.8)
fig_h = Inches(3.2)
gap   = Inches(1.5)
total = 2 * fig_w + gap
start_x = (W - total) / 2
fig_top  = Inches(1.6)

for i in range(2):
    fx = start_x + i * (fig_w + gap)
    # body shape (rounded rect standing figure outline)
    body = add_rect(slide, fx, fig_top, fig_w, fig_h,
                    fill_color=None, line_color=BORDER, line_width_pt=2.5)
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x1C)
    # head circle
    head_size = Inches(0.55)
    hx = fx + (fig_w - head_size) / 2
    hy = fig_top + Inches(0.18)
    hc = slide.shapes.add_shape(9, hx, hy, head_size, head_size)
    hc.fill.solid(); hc.fill.fore_color.rgb = BORDER; hc.line.fill.background()
    # amber "visceral fat" highlight on the right figure
    if i == 1:
        fat_w, fat_h = Inches(1.1), Inches(0.65)
        fat_x = fx + (fig_w - fat_w) / 2
        fat_y = fig_top + Inches(1.4)
        fat = slide.shapes.add_shape(9, fat_x, fat_y, fat_w, fat_h)
        fat.fill.solid(); fat.fill.fore_color.rgb = AMBER
        fat.line.color.rgb = RGBColor(0xFF, 0xD0, 0x80)
        fat.line.width = Pt(1.0)
        # tag label
        tag = add_textbox(slide, "Visceral fat", fat_x - Inches(0.5), fat_y + fat_h + Inches(0.05),
                          Inches(1.9), Inches(0.35), font_size=13, color=AMBER, align=PP_ALIGN.CENTER)
    # BMI label
    lbl_y = fig_top + fig_h + Inches(0.1)
    add_textbox(slide, "BMI 22", fx, lbl_y, fig_w, Inches(0.38),
                font_size=20, color=(TEAL if i == 0 else AMBER), align=PP_ALIGN.CENTER,
                bold=True, mono=True)
    # sub label
    add_textbox(slide, ("No hidden risk" if i == 0 else "Hidden risk"),
                fx, lbl_y + Inches(0.38), fig_w, Inches(0.38),
                font_size=18, color=(TEAL if i == 0 else RED), align=PP_ALIGN.CENTER)

# Caption
cap_y = H - MARGIN - Inches(0.55)
add_textbox(slide, "BMI measures weight. It can\u2019t see where fat is hiding.",
            MARGIN, cap_y, W - 2 * MARGIN, Inches(0.5),
            font_size=22, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

# map-pin (amber) motif slide 3
add_mappin(slide, W_IN - 1.1, 6.85, size_in=0.32, color=AMBER)


# ────────────────────────────────────────────── SLIDE 4 — The Hidden Problem ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 4)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "The Hidden Problem"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

cards = [
    ("Standard scores lose accuracy\nat normal weight",         AMBER, "~"),
    ("One major score gets it\nbackwards entirely",             RED,   "↙"),
    ("AI models that work,\nbut explain nothing",              BLUE,  "?"),
    ("Risk estimates that fail\nthe people who need them most", RED,   "⚠"),
]
card_w = (W - 2 * MARGIN - Inches(0.3)) / 2
card_h = Inches(2.0)
gap_x  = Inches(0.3)
gap_y  = Inches(0.25)
card_top = Inches(1.55)

for idx, (txt, col, icon) in enumerate(cards):
    row = idx // 2; col_i = idx % 2
    cx = MARGIN + col_i * (card_w + gap_x)
    cy = card_top + row * (card_h + gap_y)

    # card bg
    card = add_rect(slide, cx, cy, card_w, card_h, fill_color=PANEL,
                    line_color=col, line_width_pt=1.5)
    # icon char
    add_textbox(slide, icon, cx + Inches(0.2), cy + Inches(0.15),
                Inches(0.6), Inches(0.6), font_size=30, color=col, bold=True)
    # text
    tb_c = slide.shapes.add_textbox(cx + Inches(0.1), cy + Inches(0.75),
                                     card_w - Inches(0.2), card_h - Inches(0.85))
    tb_c.word_wrap = True
    tf_c = tb_c.text_frame; tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]; p_c.alignment = PP_ALIGN.LEFT
    r_c = p_c.add_run(); r_c.text = txt
    r_c.font.name = "Inter"; r_c.font.size = Pt(22); r_c.font.color.rgb = TEXT


# ─────────────────────────────────────────────────────── SLIDE 5 — The Gap ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 5)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "The Gap"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

check_items = [
    "Built specifically for normal-weight patients",
    "Risk scores tied to real medical tests",
    "Reliable confidence — even for the riskiest patients",
    "Plain explanations, not a black box",
]
item_h = Inches(0.75)
list_top = Inches(1.8)
for i, item in enumerate(check_items):
    iy = list_top + i * (item_h + Inches(0.12))
    # check box
    add_checkmark(slide, (MARGIN + Inches(0.18)) / Inches(1), iy / Inches(1) + 0.24, size_in=0.34)
    # text
    add_textbox(slide, item,
                MARGIN + Inches(0.65), iy,
                W - 2 * MARGIN - Inches(0.65), item_h,
                font_size=26, color=TEXT)

cap_y = H - MARGIN - Inches(0.6)
add_textbox(slide, "No existing system does all four. This one does.",
            MARGIN, cap_y, W - 2 * MARGIN, Inches(0.5),
            font_size=22, color=GOLD, align=PP_ALIGN.CENTER, bold=True)


# ────────────────────────────────────────────────── SLIDE 6 — How It Works ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 6)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "How It Works"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

steps = [
    ("Blood test\nvalues",         "💉", TEAL),
    ("Pattern read\nas a whole",   "🔬", BLUE),
    ("Two risk\nscores",           "📊", AMBER),
    ("Patient placed\non a map",   "📍", RED),
    ("Confidence +\nnext steps",   "✅", TEAL),
]
n = len(steps)
step_w  = Inches(1.9)
step_h  = Inches(2.4)
total_w = n * step_w + (n - 1) * Inches(0.4)
sx0     = (W - total_w) / 2
sy      = Inches(2.0)
arr_y   = sy + step_h / 2

for i, (label, icon, col) in enumerate(steps):
    sx = sx0 + i * (step_w + Inches(0.4))
    # bubble
    bubble = add_rect(slide, sx, sy, step_w, step_h,
                      fill_color=PANEL, line_color=col, line_width_pt=2.0)
    # icon label
    add_textbox(slide, icon, sx, sy + Inches(0.25),
                step_w, Inches(0.7), font_size=30, align=PP_ALIGN.CENTER)
    # text
    tb_s = slide.shapes.add_textbox(sx + Inches(0.1), sy + Inches(1.0),
                                     step_w - Inches(0.2), step_h - Inches(1.15))
    tb_s.word_wrap = True
    tf_s = tb_s.text_frame; tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]; p_s.alignment = PP_ALIGN.CENTER
    r_s = p_s.add_run(); r_s.text = label
    r_s.font.name = "Inter"; r_s.font.size = Pt(21); r_s.font.color.rgb = TEXT

    # arrow between steps
    if i < n - 1:
        arx = sx + step_w + Inches(0.02)
        arr_box = add_rect(slide, arx, arr_y - Inches(0.05),
                           Inches(0.38), Inches(0.1), fill_color=BORDER)
        arr_box.line.fill.background()
        add_textbox(slide, "→", arx, arr_y - Inches(0.28),
                    Inches(0.38), Inches(0.45), font_size=24, color=MUTED, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────── SLIDE 7 — The Map ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 7)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "The Map"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

map_size = Inches(4.0)
mx = (W - map_size) / 2
my = Inches(1.55)

# background for map
add_rect(slide, mx, my, map_size, map_size, fill_color=PANEL,
         line_color=BORDER, line_width_pt=2.0)

# four quadrants (2x2)
qw, qh = map_size / 2, map_size / 2
quads = [
    (0, 0, BLUE,  "Liver\nRisk"),
    (1, 0, RED,   "Dual\nBurden"),
    (0, 1, TEAL,  "Healthy"),
    (1, 1, AMBER, "IR\nRisk"),
]
for (c, r_i, col, lbl) in quads:
    qx = mx + c * qw; qy = my + r_i * qh
    # semi-transparent tinted quad
    q_rect = add_rect(slide, qx, qy, qw, qh,
                      fill_color=None, line_color=BORDER, line_width_pt=0.5)
    q_rect.fill.solid()
    q_rect.fill.fore_color.rgb = col
    # set transparency via XML
    solidFill = q_rect._element.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is None:
            srgb = solidFill.find(qn('a:schemeClr'))
        if srgb is not None:
            alpha_el = srgb.find(qn('a:alpha'))
            if alpha_el is None:
                alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', '20000')   # ~20% opacity (100000 = full)

    add_textbox(slide, lbl, qx + Inches(0.1), qy + Inches(0.1),
                qw - Inches(0.2), qh - Inches(0.2),
                font_size=18, color=col, bold=True, align=PP_ALIGN.CENTER)

# Map pin in red quadrant (top-right)
add_mappin(slide, (mx + map_size * 0.75) / Inches(1),
           (my + map_size * 0.27) / Inches(1),
           size_in=0.42, color=RED)

# Caption
cap_y = my + map_size + Inches(0.2)
add_textbox(slide, "Every patient gets a location. Location reveals hidden risk.",
            MARGIN, cap_y, W - 2 * MARGIN, Inches(0.55),
            font_size=22, color=MUTED, align=PP_ALIGN.CENTER, italic=True)


# ────────────────────────────────────────────── SLIDE 8 — Scale of Problem ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 8)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Scale of the Problem"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

# Glow panel
gw, gh = Inches(7.0), Inches(2.5)
gx = (W - gw) / 2; gy = Inches(1.9)
add_rect(slide, gx - Inches(0.2), gy - Inches(0.2),
         gw + Inches(0.4), gh + Inches(0.4),
         fill_color=PANEL, line_color=AMBER, line_width_pt=2.0)
add_rect(slide, gx, gy, gw, gh, fill_color=PANEL)

# Hero number
tb_num = slide.shapes.add_textbox(gx, gy + Inches(0.15), gw, Inches(1.6))
tf_n = tb_num.text_frame
p_n = tf_n.paragraphs[0]; p_n.alignment = PP_ALIGN.CENTER
r_n = p_n.add_run(); r_n.text = "[PREVALENCE ESTIMATE]"
r_n.font.name = "JetBrains Mono"; r_n.font.size = Pt(58)
r_n.font.bold = True; r_n.font.color.rgb = AMBER

# Sub-line
add_textbox(slide,
            "(range reflects sample size — see full report for detail)",
            gx, gy + Inches(1.85), gw, Inches(0.55),
            font_size=18, color=MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────── SLIDE 9 — Live Demonstration ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 9)

# Centered big text only
tb = slide.shapes.add_textbox(MARGIN, Inches(2.6), W - 2 * MARGIN, Inches(2.0))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Live demonstration"
r.font.name = "Inter"; r.font.size = Pt(52); r.font.bold = True; r.font.color.rgb = TEXT

# Small map-pin motif (teal)
add_mappin(slide, W_IN / 2, 5.6, size_in=0.55, color=TEAL)


# ────────────────────────────────────────── SLIDE 10 — Putting it to Test ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 10)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Putting It to the Test"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

# Horizontal bar chart data (Spearman ρ, normalized to bar widths)
chart_data = [
    ("LMSIS",      0.542, TEAL,  True),
    ("Random Forest", 0.596, MUTED, False),
    ("XGBoost",    0.581, MUTED, False),
    ("FLI",        0.399, MUTED, False),
    ("NAFLD-LFS",  -0.118, RED,  False),   # inverted
]
bar_max   = 0.65   # for scaling
chart_x   = MARGIN + Inches(1.4)
chart_top = Inches(1.55)
bar_h     = Inches(0.52)
bar_gap   = Inches(0.22)
max_bar_w = W - chart_x - MARGIN - Inches(0.8)
axis_x    = chart_x + Inches(0.02)   # zero line position

for i, (name, val, col, bold) in enumerate(chart_data):
    by = chart_top + i * (bar_h + bar_gap)
    # label
    add_textbox(slide, name, MARGIN, by + Inches(0.06),
                Inches(1.3), bar_h, font_size=17,
                color=(TEXT if bold else MUTED), align=PP_ALIGN.RIGHT, bold=bold)
    # bar
    bw = abs(val / bar_max) * max_bar_w
    if val < 0:
        # inverted bar going left from zero-ish
        add_rect(slide, axis_x - bw, by, bw, bar_h, fill_color=col)
    else:
        add_rect(slide, axis_x, by, bw, bar_h, fill_color=col)
    # value label
    val_str = f"{val:+.3f}"
    lx = (axis_x + bw + Inches(0.08)) if val >= 0 else (axis_x - bw - Inches(0.72))
    add_textbox(slide, val_str, lx, by + Inches(0.08),
                Inches(0.7), bar_h - Inches(0.1),
                font_size=15, color=col, bold=True, mono=True)

# zero line
add_rect(slide, axis_x - Pt(1), chart_top - Inches(0.1),
         Pt(2), len(chart_data) * (bar_h + bar_gap), fill_color=BORDER)

# Caption
cap_y = chart_top + len(chart_data) * (bar_h + bar_gap) + Inches(0.15)
add_textbox(slide, "One widely used method ranks sick patients as healthier.",
            MARGIN, cap_y, W - 2 * MARGIN, Inches(0.5),
            font_size=22, color=RED, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────── SLIDE 11 — Reliable for Riskiest ──────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 11)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Reliable for the Riskiest Patients"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

bar_pairs = [
    ("Before", 81.6,  MUTED),
    ("After",  90.4,  TEAL),
]
bp_w    = Inches(3.6)
bp_h    = Inches(3.2)
bp_gap  = Inches(1.4)
total_b = 2 * bp_w + bp_gap
bx0     = (W - total_b) / 2
by0     = Inches(1.6)
max_h   = Inches(2.6)

for i, (lbl, val, col) in enumerate(bar_pairs):
    bx = bx0 + i * (bp_w + bp_gap)
    bar_actual_h = max_h * val / 100
    bar_y = by0 + max_h - bar_actual_h
    # bar body
    bar_rect = add_rect(slide, bx + Inches(0.6), bar_y,
                        bp_w - Inches(1.2), bar_actual_h, fill_color=col)
    # percentage on top
    add_textbox(slide, f"{val:.1f}%",
                bx + Inches(0.4), bar_y - Inches(0.55),
                bp_w - Inches(0.8), Inches(0.48),
                font_size=32, color=col, bold=True, mono=True, align=PP_ALIGN.CENTER)
    # label below
    add_textbox(slide, lbl,
                bx, by0 + max_h + Inches(0.1), bp_w, Inches(0.45),
                font_size=24, color=(TEXT if i == 1 else MUTED), align=PP_ALIGN.CENTER, bold=(i == 1))

cap_y = by0 + max_h + Inches(0.65)
add_textbox(slide, "Standard methods miss the riskiest patients. This fixes that.",
            MARGIN, cap_y, W - 2 * MARGIN, Inches(0.5),
            font_size=22, color=MUTED, align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────── SLIDE 12 — Proven Beyond Training ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 12)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Proven Beyond the Training Data"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

stat_cards = [
    ("Tested two years later",               "[RESULT]", "🕐", TEAL),
    ("Tested on a different population group", "[RESULT]", "👥", BLUE),
]
sc_w   = (W - 2 * MARGIN - Inches(0.5)) / 2
sc_h   = Inches(3.0)
sc_top = Inches(1.9)

for i, (label, result, icon, col) in enumerate(stat_cards):
    sx = MARGIN + i * (sc_w + Inches(0.5))
    # card bg
    add_rect(slide, sx, sc_top, sc_w, sc_h, fill_color=PANEL,
             line_color=col, line_width_pt=2.0)
    # icon
    add_textbox(slide, icon, sx + Inches(0.2), sc_top + Inches(0.2),
                Inches(0.7), Inches(0.7), font_size=32, align=PP_ALIGN.LEFT)
    # label
    tb_lbl = slide.shapes.add_textbox(sx + Inches(0.2), sc_top + Inches(0.95),
                                      sc_w - Inches(0.4), Inches(0.9))
    tb_lbl.word_wrap = True
    tf_lbl = tb_lbl.text_frame; tf_lbl.word_wrap = True
    p_lbl = tf_lbl.paragraphs[0]; p_lbl.alignment = PP_ALIGN.LEFT
    r_lbl = p_lbl.add_run(); r_lbl.text = label
    r_lbl.font.name = "Inter"; r_lbl.font.size = Pt(22); r_lbl.font.color.rgb = MUTED

    # result number (big, mono)
    add_textbox(slide, result,
                sx + Inches(0.2), sc_top + Inches(1.95), sc_w - Inches(0.4), Inches(0.9),
                font_size=36, color=col, bold=True, mono=True)


# ───────────────────────────────────────────── SLIDE 13 — What We Built ──────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 13)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "What We Built"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

built_items = [
    ("Risk scores anchored to real medical tests",             TEAL,  "⚓"),
    ("Reliable confidence ranges for the riskiest patients",   BLUE,  "🛡"),
    ("Plain explanations of how the model behaves",            AMBER, "📋"),
    ("A complete, tested, working system",                     TEAL,  "✓"),
]
bi_h    = Inches(0.85)
bi_top  = Inches(1.65)
bi_gap  = Inches(0.2)

for i, (txt, col, icon) in enumerate(built_items):
    iy = bi_top + i * (bi_h + bi_gap)
    # row bg
    add_rect(slide, MARGIN, iy, W - 2 * MARGIN, bi_h, fill_color=PANEL,
             line_color=col, line_width_pt=1.5)
    # icon
    add_textbox(slide, icon, MARGIN + Inches(0.15), iy + Inches(0.1),
                Inches(0.6), bi_h - Inches(0.2), font_size=28, color=col)
    # text
    add_textbox(slide, txt,
                MARGIN + Inches(0.9), iy + Inches(0.15),
                W - 2 * MARGIN - Inches(1.1), bi_h - Inches(0.3),
                font_size=26, color=TEXT)

# map-pin motif (teal) slide 13
add_mappin(slide, W_IN - 1.0, 7.0, size_in=0.32, color=TEAL)


# ────────────────────────────────────────── SLIDE 14 — Honest About Limits ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)
add_footer_num(slide, 14)

tb = slide.shapes.add_textbox(MARGIN, MARGIN * 0.6, W - 2 * MARGIN, Inches(0.9))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Honest About the Limits"
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

# Divider line
add_rect(slide, MARGIN, Inches(4.0), W - 2 * MARGIN, Pt(1.5), fill_color=BORDER)

halves = [
    # (top_y, label, items, col)
    (Inches(1.45), "What we don\u2019t yet know",
     ["Built from a single national data source",
      "A one-time snapshot, not tracked over time",
      "Not a diagnostic device"],
     AMBER),
    (Inches(4.15), "What\u2019s next",
     ["Test on completely independent data",
      "Track real patients over time",
      "Expand to more population groups"],
     TEAL),
]
for (hy, hlabel, hitems, col) in halves:
    # half-label
    add_textbox(slide, hlabel, MARGIN, hy, W - 2 * MARGIN, Inches(0.5),
                font_size=20, color=col, bold=True)
    for j, it in enumerate(hitems):
        add_textbox(slide, f"• {it}",
                    MARGIN + Inches(0.3), hy + Inches(0.55) + j * Inches(0.62),
                    W - 2 * MARGIN - Inches(0.3), Inches(0.58),
                    font_size=23, color=(MUTED if col == AMBER else TEXT))


# ──────────────────────────────────────────────────── SLIDE 15 — Closing ──
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG)

# Horizontal teal accent line at top
add_rect(slide, MARGIN, Inches(0.6), W - 2 * MARGIN, Pt(2.5), fill_color=TEAL)

tb = slide.shapes.add_textbox(MARGIN, Inches(1.1), W - 2 * MARGIN, Inches(1.1))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Normal BMI is not metabolic health."
r.font.name = "Inter"; r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = TEXT

tb2 = slide.shapes.add_textbox(MARGIN, Inches(2.3), W - 2 * MARGIN, Inches(0.9))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "We can show you where the patient actually is \u2014 and the route back."
r2.font.name = "Inter"; r2.font.size = Pt(26); r2.font.bold = False; r2.font.color.rgb = MUTED

# Trust badges (3 small items)
badges = [("⚙ Tested", TEAL), ("♻ Reproducible", BLUE), ("⬡ Open Source", AMBER)]
badge_w = Inches(2.6)
badge_y = Inches(4.5)
bx0 = (W - 3 * badge_w - 2 * Inches(0.3)) / 2
for i, (txt, col) in enumerate(badges):
    bx = bx0 + i * (badge_w + Inches(0.3))
    add_rect(slide, bx, badge_y, badge_w, Inches(0.55),
             fill_color=PANEL, line_color=col, line_width_pt=1.5)
    add_textbox(slide, txt, bx, badge_y + Inches(0.03), badge_w, Inches(0.5),
                font_size=20, color=col, align=PP_ALIGN.CENTER)

# Final line
add_textbox(slide, "Thank you. Questions?",
            MARGIN, Inches(5.5), W - 2 * MARGIN, Inches(0.7),
            font_size=32, color=TEXT, align=PP_ALIGN.CENTER, bold=True)

# map-pin motif (red) — slide 15
add_mappin(slide, 1.1, 6.8, size_in=0.38, color=RED)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\singh\TYPE2-PHENOTYPE\LMSIS_Defense.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
